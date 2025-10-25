from langchain_text_splitters import RecursiveCharacterTextSplitter
import pdfplumber
from docx import Document as DocxDocument
import openpyxl
from typing import List, Dict
import chardet
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

class DocumentProcessor:
    SUPPORTED_FORMATS = {'.pdf', '.docx', '.txt', '.md', '.xlsx', '.xls', '.pptx'}
    
    @staticmethod
    def is_supported(file_type: str) -> bool:
        return any(file_type.lower().endswith(fmt) for fmt in DocumentProcessor.SUPPORTED_FORMATS)
    
    @staticmethod
    def extract_text(file_path: str, file_type: str) -> str:
        file_type = file_type.lower()
        
        try:
            if file_type.endswith('.pdf'):
                return DocumentProcessor._extract_pdf(file_path)
            elif file_type.endswith('.docx'):
                return DocumentProcessor._extract_docx(file_path)
            elif file_type.endswith(('.xlsx', '.xls')):
                return DocumentProcessor._extract_excel(file_path)
            elif file_type.endswith('.pptx'):
                return DocumentProcessor._extract_pptx(file_path)
            elif file_type.endswith(('.txt', '.md')):
                return DocumentProcessor._extract_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        page_text = page_text.replace('\x00', '')
                        text += f"\n--- Page {page_num + 1} ---\n"
                        text += page_text
                    
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            text += "\n[TABLE]\n"
                            for row in table:
                                text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                            text += "[/TABLE]\n"
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
            raise ValueError(f"Failed to extract PDF: {str(e)}")
        
        if not text.strip():
            raise ValueError("PDF appears to be empty or contains only images")
        
        return text.strip()
    
    @staticmethod
    def _extract_docx(file_path: str) -> str:
        try:
            doc = DocxDocument(file_path)
            text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            for table in doc.tables:
                text.append("\n[TABLE]")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    text.append(row_text)
                text.append("[/TABLE]\n")
            
            result = "\n".join(text)
            
            if not result.strip():
                raise ValueError("DOCX appears to be empty")
            
            return result.strip()
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}")
            raise ValueError(f"Failed to extract DOCX: {str(e)}")
    
    @staticmethod
    def _extract_excel(file_path: str) -> str:
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text.append(f"\n--- Sheet: {sheet_name} ---\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text.append(row_text)
            
            result = "\n".join(text)
            
            if not result.strip():
                raise ValueError("Excel file appears to be empty")
            
            return result.strip()
        except Exception as e:
            logger.error(f"Excel extraction error: {str(e)}")
            raise ValueError(f"Failed to extract Excel: {str(e)}")
    
    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides):
                text.append(f"\n--- Slide {slide_num + 1} ---\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
                    
                    if shape.has_table:
                        text.append("\n[TABLE]")
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            text.append(row_text)
                        text.append("[/TABLE]\n")
            
            result = "\n".join(text)
            
            if not result.strip():
                raise ValueError("PowerPoint appears to be empty")
            
            return result.strip()
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}")
            raise ValueError(f"Failed to extract PPTX: {str(e)}")
    
    @staticmethod
    def _extract_txt(file_path: str) -> str:
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
            
            detected = chardet.detect(raw_data)
            encoding = detected['encoding'] or 'utf-8'
            
            text = raw_data.decode(encoding, errors='replace')
            
            if not text.strip():
                raise ValueError("Text file appears to be empty")
            
            return text.strip()
        except Exception as e:
            logger.error(f"TXT extraction error: {str(e)}")
            raise ValueError(f"Failed to extract text file: {str(e)}")
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, any]]:
        if not text or not text.strip():
            raise ValueError("Cannot chunk empty text")
        
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""],
            length_function=len,
            is_separator_regex=False
        )
        
        chunks = text_splitter.split_text(text)
        
        return [
            {
                "text": chunk.strip(),
                "index": i,
                "metadata": {
                    "char_count": len(chunk),
                    "word_count": len(chunk.split()),
                }
            }
            for i, chunk in enumerate(chunks)
            if chunk.strip()
        ]